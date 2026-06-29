from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Helpers ──────────────────────────────────────────────────────────
def add_bg(slide, r, g, b):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)

def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=(255,255,255), alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(*color)
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_para(tf, text, font_size=16, bold=False, color=(220,220,220),
             font_name="Calibri", space_before=6):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(*color)
    p.font.name = font_name
    p.space_before = Pt(space_before)
    return p

def add_shape_bar(slide, left, top, width, height, r, g, b):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)
    shape.line.fill.background()

# ── Colors ───────────────────────────────────────────────────────────
BG       = (18, 18, 30)      # dark navy
ACCENT   = (0, 180, 216)     # cyan
ACCENT2  = (114, 9, 183)     # purple
WHITE    = (255, 255, 255)
LIGHT    = (200, 200, 220)
GOLD     = (255, 196, 0)
GREEN    = (0, 200, 120)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ─────────────────────────────────────────────────────────────────────
# SLIDE 1 — Title
# ─────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(sl, *BG)
add_shape_bar(sl, 0, 0, 13.333, 0.08, *ACCENT)
add_shape_bar(sl, 0, 7.42, 13.333, 0.08, *ACCENT2)

add_textbox(sl, 1, 1.5, 11, 1.2,
            "BLUESTOCK MUTUAL FUND ANALYTICS", 40, True, ACCENT, PP_ALIGN.CENTER)
add_textbox(sl, 1, 2.8, 11, 0.8,
            "DAY 2: Data Cleaning + SQL Database Design", 28, True, WHITE, PP_ALIGN.CENTER)
add_textbox(sl, 1, 4.0, 11, 0.6,
            "Capstone Project  |  Week 1", 20, False, LIGHT, PP_ALIGN.CENTER)
add_textbox(sl, 1, 5.5, 11, 0.5,
            "Prepared by: Vasu Aditya K  |  June 2026", 16, False, LIGHT, PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────
# SLIDE 2 — Objectives
# ─────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, *BG)
add_shape_bar(sl, 0, 0, 13.333, 0.06, *ACCENT)

add_textbox(sl, 0.8, 0.4, 11, 0.7, "Day 2 — Objectives", 32, True, ACCENT)

objectives = [
    "✦  Clean and validate all 10 datasets (handle nulls, duplicates, type errors)",
    "✦  Design and implement a 5+ table SQLite database schema",
    "✦  Load all cleaned data into the database",
    "✦  Write and test 10 SQL queries for basic analytics",
    "✦  Create a data dictionary documenting all columns and tables",
]
tf = add_textbox(sl, 1.0, 1.5, 10, 0.5, objectives[0], 20, False, WHITE)
for obj in objectives[1:]:
    add_para(tf, obj, 20, False, WHITE, space_before=14)

# ─────────────────────────────────────────────────────────────────────
# SLIDE 3 — Datasets Overview
# ─────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, *BG)
add_shape_bar(sl, 0, 0, 13.333, 0.06, *ACCENT)

add_textbox(sl, 0.8, 0.4, 11, 0.7, "10 Raw Datasets", 32, True, ACCENT)

datasets = [
    ("01_fund_master.csv",           "40 rows",     "Fund metadata — name, category, manager, expense ratio"),
    ("02_nav_history.csv",           "45K+ rows",   "Daily NAV values for all funds"),
    ("03_aum_by_fund_house.csv",     "90 rows",     "AUM per fund house per quarter"),
    ("04_monthly_sip_inflows.csv",   "48 rows",     "Monthly SIP inflow data"),
    ("05_category_inflows.csv",      "144 rows",    "Net inflows by fund category"),
    ("06_industry_folio_count.csv",  "21 rows",     "Industry-level folio counts"),
    ("07_scheme_performance.csv",    "40 rows",     "Returns, Sharpe, alpha, beta per scheme"),
    ("08_investor_transactions.csv", "32K+ rows",   "Investor-level transaction records"),
    ("09_portfolio_holdings.csv",    "322 rows",    "Stock holdings per fund"),
    ("10_benchmark_indices.csv",     "8K+ rows",    "NIFTY / SENSEX daily close values"),
]

y = 1.4
for fname, rows, desc in datasets:
    tf = add_textbox(sl, 1.0, y, 3.5, 0.35, fname, 14, True, GOLD)
    add_textbox(sl, 4.6, y, 1.5, 0.35, rows, 14, False, GREEN, PP_ALIGN.CENTER)
    add_textbox(sl, 6.2, y, 6.5, 0.35, desc, 14, False, LIGHT)
    y += 0.55

# ─────────────────────────────────────────────────────────────────────
# SLIDE 4 — Cleaning: NAV History
# ─────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, *BG)
add_shape_bar(sl, 0, 0, 13.333, 0.06, *ACCENT)

add_textbox(sl, 0.8, 0.4, 11, 0.7, "Task 1 — Clean nav_history.csv", 32, True, ACCENT)
add_textbox(sl, 0.8, 1.1, 5, 0.4, "Tools: Pandas, datetime", 16, False, GOLD)

steps = [
    "✔  Parsed date column to datetime format",
    "✔  Sorted by amfi_code + date",
    "✔  Forward-filled missing NAV for holidays & weekends",
    "✔  Removed duplicate (amfi_code, date) pairs",
    "✔  Validated NAV > 0 — dropped invalid rows",
    "",
    "Result: 64,320 rows  →  saved to data/processed/02_nav_history.csv",
]
tf = add_textbox(sl, 1.0, 2.0, 10, 0.5, steps[0], 20, False, WHITE)
for s in steps[1:]:
    color = GREEN if s.startswith("Result") else WHITE
    add_para(tf, s, 20, s.startswith("Result"), color, space_before=12)

# ─────────────────────────────────────────────────────────────────────
# SLIDE 5 — Cleaning: Investor Transactions
# ─────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, *BG)
add_shape_bar(sl, 0, 0, 13.333, 0.06, *ACCENT)

add_textbox(sl, 0.8, 0.4, 11, 0.7, "Task 2 — Clean investor_transactions.csv", 32, True, ACCENT)
add_textbox(sl, 0.8, 1.1, 5, 0.4, "Tools: Pandas, re", 16, False, GOLD)

steps = [
    "✔  Standardised transaction_type → SIP / Lumpsum / Redemption",
    "✔  Validated amount_inr > 0 — dropped invalid rows",
    "✔  Parsed transaction_date to YYYY-MM-DD format",
    "✔  Checked KYC status enum: Verified, Pending, Rejected, Unknown",
    "✔  Removed duplicate rows",
    "",
    "Result: 32,778 rows  →  saved to data/processed/08_investor_transactions.csv",
]
tf = add_textbox(sl, 1.0, 2.0, 11, 0.5, steps[0], 20, False, WHITE)
for s in steps[1:]:
    color = GREEN if s.startswith("Result") else WHITE
    add_para(tf, s, 20, s.startswith("Result"), color, space_before=12)

# ─────────────────────────────────────────────────────────────────────
# SLIDE 6 — Cleaning: Scheme Performance
# ─────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, *BG)
add_shape_bar(sl, 0, 0, 13.333, 0.06, *ACCENT)

add_textbox(sl, 0.8, 0.4, 11, 0.7, "Task 3 — Clean scheme_performance.csv", 32, True, ACCENT)
add_textbox(sl, 0.8, 1.1, 5, 0.4, "Tools: Pandas, NumPy", 16, False, GOLD)

steps = [
    "✔  Validated all return values are numeric (coerced errors to NaN)",
    "✔  Flagged negative Sharpe ratios (negative_sharpe_flag column)",
    "✔  Checked expense_ratio_pct range: 0.1% – 2.5%",
    "✔  Removed rows outside expense ratio bounds",
    "✔  Dropped duplicate records",
    "",
    "Result: 40 rows  →  saved to data/processed/07_scheme_performance.csv",
]
tf = add_textbox(sl, 1.0, 2.0, 11, 0.5, steps[0], 20, False, WHITE)
for s in steps[1:]:
    color = GREEN if s.startswith("Result") else WHITE
    add_para(tf, s, 20, s.startswith("Result"), color, space_before=12)

# ─────────────────────────────────────────────────────────────────────
# SLIDE 7 — Cleaning: Other 7 Datasets
# ─────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, *BG)
add_shape_bar(sl, 0, 0, 13.333, 0.06, *ACCENT)

add_textbox(sl, 0.8, 0.4, 11, 0.7, "Tasks — Clean Remaining 7 Datasets", 32, True, ACCENT)

rows_data = [
    ("01_fund_master.csv",          "40 rows",    "Dedup, date parsing"),
    ("03_aum_by_fund_house.csv",    "90 rows",    "Date parsing, dedup"),
    ("04_monthly_sip_inflows.csv",  "48 rows",    "Month format → YYYY-MM"),
    ("05_category_inflows.csv",     "144 rows",   "Month format → YYYY-MM"),
    ("06_industry_folio_count.csv", "21 rows",    "Month format → YYYY-MM"),
    ("09_portfolio_holdings.csv",   "322 rows",   "Date parsing, dedup"),
    ("10_benchmark_indices.csv",    "8,050 rows", "Date parsing, dedup"),
]

add_textbox(sl, 1.0, 1.3, 4.5, 0.4, "FILE", 16, True, ACCENT)
add_textbox(sl, 5.5, 1.3, 2.0, 0.4, "ROWS", 16, True, ACCENT, PP_ALIGN.CENTER)
add_textbox(sl, 7.5, 1.3, 4.5, 0.4, "CLEANING STEPS", 16, True, ACCENT)
add_shape_bar(sl, 1.0, 1.72, 10.5, 0.02, *ACCENT)

y = 1.9
for fname, rows, cleaning in rows_data:
    add_textbox(sl, 1.0, y, 4.5, 0.35, fname, 16, False, GOLD)
    add_textbox(sl, 5.5, y, 2.0, 0.35, rows, 16, False, GREEN, PP_ALIGN.CENTER)
    add_textbox(sl, 7.5, y, 4.5, 0.35, cleaning, 16, False, LIGHT)
    y += 0.55

tf = add_textbox(sl, 1.0, y + 0.4, 10, 0.5,
                 "All 10 cleaned CSVs saved to  data/processed/", 20, True, GREEN)

# ─────────────────────────────────────────────────────────────────────
# SLIDE 8 — SQLite Star Schema
# ─────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, *BG)
add_shape_bar(sl, 0, 0, 13.333, 0.06, *ACCENT)

add_textbox(sl, 0.8, 0.4, 11, 0.7, "Task 4 — SQLite Star Schema Design", 32, True, ACCENT)
add_textbox(sl, 0.8, 1.1, 5, 0.4, "Tools: SQLite, SQL DDL", 16, False, GOLD)

tables = [
    ("dim_fund",              "PK: amfi_code",                 "Fund master dimension — name, category, manager"),
    ("fact_nav",              "PK: amfi_code + date",          "Daily NAV values (FK → dim_fund)"),
    ("fact_transactions",     "FK: amfi_code",                 "Investor transaction records (FK → dim_fund)"),
    ("fact_performance",      "PK: amfi_code",                 "Scheme returns, risk metrics (FK → dim_fund)"),
    ("fact_aum",              "—",                             "AUM by fund house per quarter"),
    ("monthly_sip_inflows",   "—",                             "Monthly SIP data"),
    ("category_inflows",      "—",                             "Inflows by fund category"),
    ("industry_folio_count",  "—",                             "Folio counts by type"),
    ("portfolio_holdings",    "FK: amfi_code",                 "Stock holdings per fund"),
    ("benchmark_indices",     "—",                             "NIFTY / SENSEX daily values"),
]

add_textbox(sl, 1.0, 1.8, 3.0, 0.4, "TABLE", 14, True, ACCENT)
add_textbox(sl, 4.2, 1.8, 3.0, 0.4, "KEYS", 14, True, ACCENT, PP_ALIGN.CENTER)
add_textbox(sl, 7.4, 1.8, 5.0, 0.4, "DESCRIPTION", 14, True, ACCENT)
add_shape_bar(sl, 1.0, 2.2, 10.5, 0.02, *ACCENT)

y = 2.35
for tbl, keys, desc in tables:
    add_textbox(sl, 1.0, y, 3.0, 0.3, tbl, 14, True, GOLD)
    add_textbox(sl, 4.2, y, 3.0, 0.3, keys, 13, False, LIGHT, PP_ALIGN.CENTER)
    add_textbox(sl, 7.4, y, 5.0, 0.3, desc, 13, False, LIGHT)
    y += 0.47

# ─────────────────────────────────────────────────────────────────────
# SLIDE 9 — Data Loading
# ─────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, *BG)
add_shape_bar(sl, 0, 0, 13.333, 0.06, *ACCENT)

add_textbox(sl, 0.8, 0.4, 11, 0.7, "Task 5 — Load Data into SQLite", 32, True, ACCENT)
add_textbox(sl, 0.8, 1.1, 5, 0.4, "Tools: SQLAlchemy, SQLite", 16, False, GOLD)

steps = [
    "✔  Used SQLAlchemy create_engine('sqlite:///bluestock_mf.db')",
    "✔  Loaded each cleaned CSV via  df.to_sql(table, engine)",
    "✔  Schema enforced via schema.sql (CREATE TABLE statements)",
    "✔  Old database auto-deleted on re-run to prevent duplicates",
    "",
    "Row Count Verification:",
    "    dim_fund: 40  |  fact_nav: 64,320  |  fact_transactions: 32,778",
    "    fact_performance: 40  |  fact_aum: 90  |  benchmark_indices: 8,050",
    "    portfolio_holdings: 322  |  monthly_sip_inflows: 48",
    "",
    "Output: bluestock_mf.db  ✓",
]
tf = add_textbox(sl, 1.0, 2.0, 11, 0.5, steps[0], 18, False, WHITE)
for s in steps[1:]:
    bold = s.startswith("Row Count") or s.startswith("Output")
    color = GREEN if s.startswith("Output") else GOLD if s.startswith("Row Count") else WHITE
    add_para(tf, s, 18, bold, color, space_before=8)

# ─────────────────────────────────────────────────────────────────────
# SLIDE 10 — 10 SQL Queries
# ─────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, *BG)
add_shape_bar(sl, 0, 0, 13.333, 0.06, *ACCENT)

add_textbox(sl, 0.8, 0.4, 11, 0.7, "Task 6 — 10 Analytical SQL Queries", 32, True, ACCENT)
add_textbox(sl, 0.8, 1.1, 5, 0.4, "File: sql/queries.sql", 16, False, GOLD)

queries = [
    "Q1.  Top 5 funds by AUM",
    "Q2.  Average NAV per month for each fund",
    "Q3.  SIP inflow Year-over-Year (YoY) growth",
    "Q4.  Total transactions by state",
    "Q5.  Funds with expense ratio < 1%",
    "Q6.  Total Lumpsum amount invested per fund",
    "Q7.  Top 10 stocks by portfolio weight across all funds",
    "Q8.  Best performing funds by 3-year return",
    "Q9.  Monthly trend of total SIP accounts",
    "Q10. Average transaction amount by age group & gender",
]
tf = add_textbox(sl, 1.0, 2.0, 10, 0.5, queries[0], 20, False, WHITE)
for q in queries[1:]:
    add_para(tf, q, 20, False, WHITE, space_before=10)

# ─────────────────────────────────────────────────────────────────────
# SLIDE 11 — Sample Query Results
# ─────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, *BG)
add_shape_bar(sl, 0, 0, 13.333, 0.06, *ACCENT)

add_textbox(sl, 0.8, 0.4, 11, 0.7, "Sample Query Results", 32, True, ACCENT)

# Q1 result
add_textbox(sl, 0.8, 1.3, 6, 0.5, "Q1: Top 5 Funds by AUM (₹ Crore)", 20, True, GOLD)
q1 = [
    "1.  Mirae Asset Emerging Bluechip    — ₹49,046 Cr",
    "2.  Kotak Emerging Equity Fund       — ₹47,469 Cr",
    "3.  Nippon India Small Cap Fund      — ₹43,630 Cr",
    "4.  DSP Top 100 Equity Fund          — ₹41,828 Cr",
    "5.  UTI Mid Cap Fund                 — ₹41,728 Cr",
]
tf = add_textbox(sl, 1.0, 2.0, 6, 0.4, q1[0], 16, False, WHITE)
for q in q1[1:]:
    add_para(tf, q, 16, False, WHITE, space_before=6)

# Q5 result
add_textbox(sl, 7.0, 1.3, 6, 0.5, "Q5: Funds with Expense Ratio < 1%", 20, True, GOLD)
q5 = [
    "1.  Nippon India Gilt Securities  — 0.55%",
    "2.  HDFC Short Term Debt Fund     — 0.56%",
    "3.  Kotak Liquid Fund             — 0.60%",
    "4.  SBI Bluechip Direct           — 0.66%",
    "5.  SBI Small Cap Direct          — 0.72%",
]
tf = add_textbox(sl, 7.2, 2.0, 5.5, 0.4, q5[0], 16, False, WHITE)
for q in q5[1:]:
    add_para(tf, q, 16, False, WHITE, space_before=6)

# Q4 result
add_textbox(sl, 0.8, 4.5, 6, 0.5, "Q4: Top 5 States by Transaction Amount", 20, True, GOLD)
q4 = [
    "1.  Punjab         — ₹31.6 Cr",
    "2.  Tamil Nadu     — ₹31.5 Cr",
    "3.  Madhya Pradesh — ₹30.8 Cr",
    "4.  Rajasthan      — ₹29.9 Cr",
    "5.  Gujarat        — ₹29.8 Cr",
]
tf = add_textbox(sl, 1.0, 5.1, 5.5, 0.4, q4[0], 16, False, WHITE)
for q in q4[1:]:
    add_para(tf, q, 16, False, WHITE, space_before=6)

# ─────────────────────────────────────────────────────────────────────
# SLIDE 12 — Data Dictionary
# ─────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, *BG)
add_shape_bar(sl, 0, 0, 13.333, 0.06, *ACCENT)

add_textbox(sl, 0.8, 0.4, 11, 0.7, "Task 7 — Data Dictionary", 32, True, ACCENT)
add_textbox(sl, 0.8, 1.1, 5, 0.4, "File: data_dictionary.md", 16, False, GOLD)

items = [
    "✔  Documented all 10 tables with column-level detail",
    "✔  Included data types (TEXT, REAL, INTEGER, DATE)",
    "✔  Defined primary keys and foreign key relationships",
    "✔  Added business definitions for every column",
    "✔  Referenced source CSV for each table",
    "✔  Noted cleaning transformations applied",
]
tf = add_textbox(sl, 1.0, 2.0, 10, 0.5, items[0], 20, False, WHITE)
for item in items[1:]:
    add_para(tf, item, 20, False, WHITE, space_before=14)

# ─────────────────────────────────────────────────────────────────────
# SLIDE 13 — Deliverables Summary
# ─────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, *BG)
add_shape_bar(sl, 0, 0, 13.333, 0.06, *ACCENT)

add_textbox(sl, 0.8, 0.4, 11, 0.7, "Day 2 — Deliverables", 32, True, ACCENT)

deliverables = [
    ("data/processed/",      "10 cleaned CSV files"),
    ("bluestock_mf.db",      "SQLite database with 10 tables loaded"),
    ("sql/schema.sql",       "CREATE TABLE statements with PK / FK"),
    ("sql/queries.sql",      "10 analytical SQL queries"),
    ("data_dictionary.md",   "Column docs, types, business definitions"),
    ("src/data_cleaning.py", "Python script for data cleaning"),
    ("src/load_to_sqlite.py","Python script for DB loading"),
]

add_textbox(sl, 1.5, 1.5, 4.0, 0.4, "FILE / FOLDER", 16, True, ACCENT)
add_textbox(sl, 6.0, 1.5, 6.0, 0.4, "DESCRIPTION", 16, True, ACCENT)
add_shape_bar(sl, 1.5, 1.9, 9.5, 0.02, *ACCENT)

y = 2.1
for path, desc in deliverables:
    add_textbox(sl, 1.5, y, 4.0, 0.4, path, 18, True, GOLD)
    add_textbox(sl, 6.0, y, 6.0, 0.4, desc, 18, False, LIGHT)
    y += 0.55

tf = add_textbox(sl, 1.5, y + 0.5, 9, 0.5,
                 'Git Commit: "Day 2: Cleaned data + SQLite DB loaded"  ✓', 20, True, GREEN)

# ─────────────────────────────────────────────────────────────────────
# SLIDE 14 — Thank You
# ─────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, *BG)
add_shape_bar(sl, 0, 0, 13.333, 0.08, *ACCENT)
add_shape_bar(sl, 0, 7.42, 13.333, 0.08, *ACCENT2)

add_textbox(sl, 1, 2.5, 11, 1.0,
            "Thank You!", 44, True, ACCENT, PP_ALIGN.CENTER)
add_textbox(sl, 1, 3.8, 11, 0.6,
            "Day 2 Complete — Data Cleaned & Database Ready", 24, False, WHITE, PP_ALIGN.CENTER)
add_textbox(sl, 1, 5.0, 11, 0.5,
            "Next → Day 3: EDA + Visualization", 20, False, LIGHT, PP_ALIGN.CENTER)

# ── Save ─────────────────────────────────────────────────────────────
output_path = r'reports\Day2_Data_Cleaning_SQL_Database.pptx'
import os
os.makedirs('reports', exist_ok=True)
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
